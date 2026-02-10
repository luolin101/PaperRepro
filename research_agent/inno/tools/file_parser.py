"""
File parser module that supports reading various file formats including:
text files, PDF, DOCX, JSON, XML, YAML, HTML, LaTeX, and more.
"""
import os
import json
from abc import ABC, abstractmethod
from typing import BinaryIO
from pathlib import Path

# Optional imports - gracefully handle missing dependencies
try:
    import charset_normalizer
except ImportError:
    charset_normalizer = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import yaml
except ImportError:
    yaml = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    from pylatexenc.latex2text import LatexNodes2Text
except ImportError:
    LatexNodes2Text = None

try:
    import textract
except ImportError:
    textract = None

try:
    import docx2txt
except ImportError:
    docx2txt = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None


class ParserStrategy(ABC):
    """Abstract base class for file parsers."""
    
    @abstractmethod
    def read(self, file: BinaryIO) -> str:
        """Read and parse file content as text."""
        pass


class TXTParser(ParserStrategy):
    """Parser for plain text files."""
    
    def read(self, file: BinaryIO) -> str:
        if charset_normalizer:
            charset_match = charset_normalizer.from_bytes(file.read()).best()
            if charset_match:
                return str(charset_match)
        # Fallback to UTF-8 with error handling
        file.seek(0)
        try:
            return file.read().decode('utf-8')
        except UnicodeDecodeError:
            file.seek(0)
            return file.read().decode('utf-8', errors='replace')


class PDFParser(ParserStrategy):
    """Parser for PDF files."""
    
    def read(self, file: BinaryIO) -> str:
        if not pypdf:
            raise ValueError("pypdf library is required to read PDF files. Install it with: pip install pypdf")
        parser = pypdf.PdfReader(file)
        text = ""
        for page in parser.pages:
            text += page.extract_text() + "\n"
        return text


class DOCXParser(ParserStrategy):
    """Parser for DOCX files."""
    
    def read(self, file: BinaryIO) -> str:
        if not docx:
            raise ValueError("python-docx library is required to read DOCX files. Install it with: pip install python-docx")
        doc_file = docx.Document(file)
        text = ""
        for para in doc_file.paragraphs:
            text += para.text + "\n"
        return text


class DOCParser(ParserStrategy):
    """Parser for DOC files (old Microsoft Word binary format)."""
    
    def read(self, file: BinaryIO) -> str:
        # textract and docx2txt require file path, not file handle
        # Save content to temp file first
        file.seek(0)
        file_content = file.read()
        
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix='.doc') as tmp_file:
            tmp_file.write(file_content)
            tmp_path = tmp_file.name
        
        try:
            # Try textract first (most comprehensive)
            if textract:
                try:
                    text = textract.process(tmp_path).decode('utf-8', errors='replace')
                    return text
                except Exception as e:
                    # If textract fails, try fallback
                    pass
            
            # Fallback: Try docx2txt (might work for some .doc files)
            if docx2txt:
                try:
                    text = docx2txt.process(tmp_path)
                    if text:
                        return text
                except Exception:
                    pass
            
            # If no library available or all methods failed, provide helpful error message
            raise ValueError(
                "textract library is required to read DOC files. "
                "Install it with: pip install textract\n"
                "Note: textract may require additional system dependencies. "
                "Alternatively, you can convert .doc files to .docx format."
            )
        finally:
            # Always clean up temp file
            try:
                os.unlink(tmp_path)
            except Exception:
                pass


class JSONParser(ParserStrategy):
    """Parser for JSON files."""
    
    def read(self, file: BinaryIO) -> str:
        data = json.load(file)
        return json.dumps(data, indent=2, ensure_ascii=False)


class XMLParser(ParserStrategy):
    """Parser for XML files."""
    
    def read(self, file: BinaryIO) -> str:
        if not BeautifulSoup:
            raise ValueError("beautifulsoup4 library is required to read XML files. Install it with: pip install beautifulsoup4")
        soup = BeautifulSoup(file, "xml")
        return soup.get_text()


class YAMLParser(ParserStrategy):
    """Parser for YAML files."""
    
    def read(self, file: BinaryIO) -> str:
        if not yaml:
            raise ValueError("pyyaml library is required to read YAML files. Install it with: pip install pyyaml")
        data = yaml.safe_load(file)
        return str(data)


class HTMLParser(ParserStrategy):
    """Parser for HTML files."""
    
    def read(self, file: BinaryIO) -> str:
        if not BeautifulSoup:
            raise ValueError("beautifulsoup4 library is required to read HTML files. Install it with: pip install beautifulsoup4")
        soup = BeautifulSoup(file, "html.parser")
        return soup.get_text()


class LaTeXParser(ParserStrategy):
    """Parser for LaTeX files."""
    
    def read(self, file: BinaryIO) -> str:
        if not LatexNodes2Text:
            raise ValueError("pylatexenc library is required to read LaTeX files. Install it with: pip install pylatexenc")
        latex = file.read().decode('utf-8', errors='replace')
        return LatexNodes2Text().latex_to_text(latex)


class RTFParser(ParserStrategy):
    """Parser for RTF files - basic implementation."""
    
    def read(self, file: BinaryIO) -> str:
        # RTF is complex, this is a simple text extraction
        content = file.read().decode('latin-1', errors='replace')
        # Remove RTF control words (simplified)
        import re
        # Remove RTF control sequences
        content = re.sub(r'\\[a-z]+\d*\s?', '', content)
        content = re.sub(r'[{}\\]', '', content)
        return content


class XLSXParser(ParserStrategy):
    """Parser for Excel files (.xlsx)."""
    
    def read(self, file: BinaryIO) -> str:
        if not pd:
            raise ValueError("pandas library is required to read Excel files. Install it with: pip install pandas openpyxl")
        # pandas requires file path or file-like object
        file.seek(0)
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_file:
            tmp_file.write(file.read())
            tmp_path = tmp_file.name
        
        try:
            sheets = pd.read_excel(tmp_path, sheet_name=None, engine='openpyxl')
            text = ""
            for sheet_name, df in sheets.items():
                text += f"=== Sheet: {sheet_name} ===\n"
                text += df.to_string(index=False) + "\n\n"
            return text
        finally:
            try:
                os.unlink(tmp_path)
            except Exception:
                pass


class XLSParser(ParserStrategy):
    """Parser for old Excel files (.xls)."""
    
    def read(self, file: BinaryIO) -> str:
        if not pd:
            raise ValueError("pandas library is required to read Excel files. Install it with: pip install pandas xlrd")
        file.seek(0)
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xls') as tmp_file:
            tmp_file.write(file.read())
            tmp_path = tmp_file.name
        
        try:
            sheets = pd.read_excel(tmp_path, sheet_name=None, engine='xlrd')
            text = ""
            for sheet_name, df in sheets.items():
                text += f"=== Sheet: {sheet_name} ===\n"
                text += df.to_string(index=False) + "\n\n"
            return text
        finally:
            try:
                os.unlink(tmp_path)
            except Exception:
                pass


class PPTXParser(ParserStrategy):
    """Parser for PowerPoint files (.pptx)."""
    
    def read(self, file: BinaryIO) -> str:
        if not pptx:
            raise ValueError("python-pptx library is required to read PowerPoint files. Install it with: pip install python-pptx")
        file.seek(0)
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
            tmp_file.write(file.read())
            tmp_path = tmp_file.name
        
        try:
            prs = pptx.Presentation(tmp_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"=== Slide {slide_num} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                text += "\n"
            return text
        finally:
            try:
                os.unlink(tmp_path)
            except Exception:
                pass


class PPTParser(ParserStrategy):
    """Parser for old PowerPoint files (.ppt)."""
    
    def read(self, file: BinaryIO) -> str:
        # .ppt is more complex, use textract as fallback
        file.seek(0)
        file_content = file.read()
        
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix='.ppt') as tmp_file:
            tmp_file.write(file_content)
            tmp_path = tmp_file.name
        
        try:
            if textract:
                try:
                    text = textract.process(tmp_path).decode('utf-8', errors='replace')
                    return text
                except Exception:
                    pass
            raise ValueError(
                "textract library is required to read PPT files. "
                "Install it with: pip install textract"
            )
        finally:
            try:
                os.unlink(tmp_path)
            except Exception:
                pass


class TOMLParser(ParserStrategy):
    """Parser for TOML files."""
    
    def read(self, file: BinaryIO) -> str:
        file.seek(0)
        content = file.read().decode('utf-8', errors='replace')
        try:
            import tomllib  # Python 3.11+
            data = tomllib.loads(content)
            return json.dumps(data, indent=2, ensure_ascii=False)
        except ImportError:
            try:
                import tomli  # Fallback for older Python
                data = tomli.loads(content)
                return json.dumps(data, indent=2, ensure_ascii=False)
            except ImportError:
                # Fallback to text parser if no TOML library available
                return content


class INIParser(ParserStrategy):
    """Parser for INI/config files."""
    
    def read(self, file: BinaryIO) -> str:
        import configparser
        config = configparser.ConfigParser()
        file.seek(0)
        content = file.read().decode('utf-8', errors='replace')
        try:
            config.read_string(content)
        except configparser.Error:
            # If parsing fails, return as plain text
            return content
        
        # Convert to readable text format
        text = ""
        for section in config.sections():
            text += f"[{section}]\n"
            for key, value in config.items(section):
                text += f"{key} = {value}\n"
            text += "\n"
        
        # Also include items not in sections (if any)
        if config.defaults():
            text += "[DEFAULT]\n"
            for key, value in config.defaults().items():
                text += f"{key} = {value}\n"
            text += "\n"
        
        return text if text else content


# Extension to parser mapping
extension_to_parser = {
    # Text files
    ".txt": TXTParser(),
    ".md": TXTParser(),
    ".markdown": TXTParser(),
    ".csv": TXTParser(),
    ".log": TXTParser(),
    # Code files (treated as text)
    ".py": TXTParser(),
    ".R": TXTParser(),
    ".r": TXTParser(),
    ".do": TXTParser(),
    ".sh": TXTParser(),
    ".bash": TXTParser(),
    ".zsh": TXTParser(),
    ".js": TXTParser(),
    ".ts": TXTParser(),
    ".tsx": TXTParser(),
    ".jsx": TXTParser(),
    ".java": TXTParser(),
    ".c": TXTParser(),
    ".cpp": TXTParser(),
    ".cc": TXTParser(),
    ".h": TXTParser(),
    ".hpp": TXTParser(),
    ".go": TXTParser(),
    ".rs": TXTParser(),
    ".php": TXTParser(),
    ".rb": TXTParser(),
    ".swift": TXTParser(),
    ".kt": TXTParser(),
    ".scala": TXTParser(),
    # Documents
    ".pdf": PDFParser(),
    ".docx": DOCXParser(),
    ".doc": DOCParser(),
    ".rtf": RTFParser(),
    ".tex": LaTeXParser(),
    # Data formats
    ".json": JSONParser(),
    ".xml": XMLParser(),
    ".yaml": YAMLParser(),
    ".yml": YAMLParser(),
    ".toml": TOMLParser(),
    # Web formats
    ".html": HTMLParser(),
    ".htm": HTMLParser(),
    ".xhtml": HTMLParser(),
    # Spreadsheets
    ".xlsx": XLSXParser(),
    ".xls": XLSParser(),
    # Presentations
    ".pptx": PPTXParser(),
    ".ppt": PPTParser(),
    # Config files
    ".ini": INIParser(),
    ".cfg": INIParser(),
    ".conf": INIParser(),
    ".properties": INIParser(),
}


def is_file_binary(file_path: str) -> bool:
    """Check if a file appears to be binary."""
    try:
        with open(file_path, 'rb') as f:
            chunk = f.read(8192)
            if b'\x00' in chunk:
                return True
        return False
    except Exception:
        return False


def decode_textual_file(file_path: str) -> str:
    """
    Read and decode a textual file based on its extension.
    
    Args:
        file_path: Path to the file to read.
    
    Returns:
        Decoded text content of the file.
    
    Raises:
        FileNotFoundError: If the file doesn't exist.
        ValueError: If the file format is not supported.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    file_ext = Path(file_path).suffix.lower()
    parser = extension_to_parser.get(file_ext)
    
    if not parser:
        # Check if file is binary - if so, we can't parse it
        if is_file_binary(file_path):
            raise ValueError(f"Unsupported binary file format: {file_ext}. Supported formats: {', '.join(extension_to_parser.keys())}")
        # Fallback to text parser for unknown extensions (likely code files)
        parser = TXTParser()
    
    # Read file in binary mode and parse
    with open(file_path, 'rb') as f:
        return parser.read(f)


def test_file_parser(file_path: str) -> dict:
    """
    Test function to verify file parsing works correctly.
    
    Args:
        file_path: Path to the file to test.
    
    Returns:
        Dictionary with test results including:
        - success: bool indicating if parsing was successful
        - file_path: the tested file path
        - extension: detected file extension
        - content_length: length of parsed content
        - content_preview: first 200 characters of parsed content
        - error: error message if parsing failed
    """
    result = {
        "success": False,
        "file_path": file_path,
        "extension": None,
        "content_length": 0,
        "content_preview": "",
        "error": None
    }
    
    try:
        if not os.path.exists(file_path):
            result["error"] = f"File not found: {file_path}"
            return result
        
        file_ext = Path(file_path).suffix.lower()
        result["extension"] = file_ext
        
        # Try to parse the file
        content = decode_textual_file(file_path)
        
        result["success"] = True
        result["content_length"] = len(content)
        result["content_preview"] = content + "..." if len(content) > 200 else content
        
        print(f"✓ Successfully parsed {file_path}")
        print(f"  Extension: {file_ext}")
        print(f"  Content length: {len(content)} characters")
        print(f"  Preview: {result['content_preview']}")
        
    except Exception as e:
        result["error"] = str(e)
        result["success"] = False
        print(f"✗ Failed to parse {file_path}")
        print(f"  Error: {str(e)}")
    
    return result


if __name__ == "__main__":
    """
    Simple test script. Usage:
    python -m research_agent.inno.tools.file_parser <file_path>
    """
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python file_parser.py <file_path>")
        print("\nExample:")
        print("  python file_parser.py /path/to/file.txt")
        print("  python file_parser.py /path/to/file.pdf")
        print("  python file_parser.py /path/to/file.docx")
        sys.exit(1)
    
    file_path = sys.argv[1]
    result = test_file_parser(file_path)
    
    # Print summary
    print("\n" + "="*60)
    print("Test Summary:")
    print("="*60)
    for key, value in result.items():
        if key == "content_preview" and value:
            print(f"{key}: {value[:100]}..." if len(str(value)) > 100 else f"{key}: {value}")
        else:
            print(f"{key}: {value}")

